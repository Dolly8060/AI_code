import os
import math
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from playwright.sync_api import sync_playwright


SLIDE_W_IN = 13.333  # 16:9 widescreen (inches)
SLIDE_H_IN = 7.5

# 分辨率预设 (width, height, scale)
RESOLUTIONS = {
    '1080p': (1920, 1080, 1),
    '2k':    (1920, 1080, 2),   # 3840x2160 
    '4k':    (1920, 1080, 2),   # 实际输出3840x2160，最清晰
}

# 默认使用4K（2x DPI）
RENDER_W_PX = 1920
RENDER_H_PX = 1080
DEVICE_SCALE = 2  # 提升清晰度关键！


def is_url(s: str) -> bool:
    return s.startswith("http://") or s.startswith("https://")


def ensure_dir(p: Path):
    p.mkdir(parents=True, exist_ok=True)


def html_to_slide_images(src: str, out_dir: Path, bg_color: str = "white", scale: int = 2):
    """
    将 HTML 页面按 16:9 的 viewport 分段截图，输出为 out_dir/slide_001.png ...
    返回图片路径列表。
    
    scale: 1=1080p, 2=4K(推荐)
    """
    ensure_dir(out_dir)
    images = []
    
    actual_w = RENDER_W_PX * scale
    actual_h = RENDER_H_PX * scale
    print(f"渲染分辨率: {actual_w} × {actual_h} (scale={scale}x)")

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(
            viewport={"width": RENDER_W_PX, "height": RENDER_H_PX},
            device_scale_factor=scale,  # 提升清晰度
        )

        # 保证背景一致 + 隐藏滚动条
        page.add_style_tag(content=f"""
            html, body {{ background: {bg_color} !important; }}
            * {{ scrollbar-width: none !important; -ms-overflow-style: none !important; }}
            *::-webkit-scrollbar {{ display: none !important; }}
        """)

        if is_url(src):
            page.goto(src, wait_until="networkidle")
        else:
            html_path = Path(src).resolve()
            if not html_path.exists():
                raise FileNotFoundError(f"找不到文件: {html_path}")
            page.goto(html_path.as_uri(), wait_until="networkidle")

        # 等待字体加载（避免字体晚到导致截图错位）
        page.wait_for_timeout(500)
        page.evaluate("() => document.fonts && document.fonts.ready")

        total_height = page.evaluate("() => document.documentElement.scrollHeight")
        total_width = page.evaluate("() => document.documentElement.scrollWidth")

        # 若页面宽度大于 viewport，优先用页面实际宽度（否则会横向缩放/换行变化）
        # 这里为了“1:1 视觉”，建议强制把内容适配到 1920 宽；
        # 如果你的 HTML 本来就是 1920 宽的设计稿，此处不需要改。
        # 如需自动适配，可在这里注入 CSS 做缩放，但会改变布局。
        # print(f"scrollWidth={total_width}, scrollHeight={total_height}")

        num_slides = math.ceil(total_height / RENDER_H_PX)

        for i in range(num_slides):
            y = i * RENDER_H_PX
            clip_h = min(RENDER_H_PX, total_height - y)

            img_path = out_dir / f"slide_{i+1:03d}.png"
            page.screenshot(
                path=str(img_path),
                full_page=False,
                clip={"x": 0, "y": y, "width": RENDER_W_PX, "height": clip_h},
            )

            # 如果最后一页高度不足 1080，补白到 1080（保证铺满 PPT 不拉伸）
            if clip_h != RENDER_H_PX:
                im = Image.open(img_path).convert("RGBA")
                canvas = Image.new("RGBA", (RENDER_W_PX, RENDER_H_PX), (255, 255, 255, 255))
                canvas.paste(im, (0, 0))
                canvas.convert("RGB").save(img_path)

            images.append(img_path)

        browser.close()

    return images


def images_to_pptx(images, pptx_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    blank_layout = prs.slide_layouts[6]  # 空白页

    for img in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img),
            left=Inches(0),
            top=Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(str(pptx_path))


def main():
    parser = argparse.ArgumentParser(description="HTML -> PPTX (pixel-perfect via screenshots)")
    parser.add_argument("src", help="HTML 本地路径 或 URL（例如 ai_zone_v2.html 或 https://...）")
    parser.add_argument("-o", "--out", default="output.pptx", help="输出 pptx 文件名")
    parser.add_argument("--tmp", default="_slides_tmp", help="中间截图目录")
    parser.add_argument("--bg", default="white", help="背景色（默认 white）")
    parser.add_argument("-s", "--scale", type=int, default=2, choices=[1, 2, 3, 4],
                        help="清晰度: 1=1080p, 2=4K(推荐), 3=5K, 4=8K")
    args = parser.parse_args()

    tmp_dir = Path(args.tmp)
    pptx_path = Path(args.out)

    imgs = html_to_slide_images(args.src, tmp_dir, bg_color=args.bg, scale=args.scale)
    images_to_pptx(imgs, pptx_path)

    print(f"完成：{pptx_path.resolve()}")
    print(f"生成页数：{len(imgs)}（中间图片在 {tmp_dir.resolve()}）")


if __name__ == "__main__":
    main()
"""CSNDC PPT Converter - 4K高清HTML转PPT工具

By Dolly

功能：
- 智能识别各类HTML slide结构
- 4K高清输出 (4000x2252)
- 使用系统Edge/Chrome渲染

打包: pyinstaller --onefile --windowed --name "CSNDC_PPT_Converter" html2ppt_app.py
"""

import os
import sys
import asyncio
import tempfile
import shutil
import threading
import re
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pathlib import Path
from typing import List, Tuple, Optional

# Windows事件循环策略：使用ProactorEventLoop支持子进程
# 注意：Playwright需要子进程支持，不能用SelectorEventLoop
if sys.platform == 'win32':
    asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())


def check_dependencies():
    """检查依赖"""
    missing = []
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        missing.append("beautifulsoup4")
    try:
        from pptx import Presentation
    except ImportError:
        missing.append("python-pptx")
    try:
        from playwright.async_api import async_playwright
    except ImportError:
        missing.append("playwright")
    try:
        from PIL import Image
    except ImportError:
        missing.append("Pillow")
    return missing


class HTMLToPPT:
    """通用HTML转PPT转换器（增强版）
    
    支持多种HTML幻灯片结构：
    - Reveal.js / Impress.js / Slidev 等框架
    - 自定义slide/page容器
    - 单页面整页转换
    """
    
    VIEWPORT_W = 1000
    VIEWPORT_H = 563
    SCALE = 4  # 4K输出
    
    # 扩展的slide选择器（优先级从高到低）
    SLIDE_SELECTORS = [
        # 常见slide框架
        ('section', 'slide'),
        ('div', 'slide'),
        ('article', 'slide'),
        # Reveal.js
        ('section', 'present'),
        ('section', None),  # Reveal.js默认用section
        # 自定义页面类
        ('div', 'page'),
        ('div', 'swiper-slide'),
        ('div', 'carousel-item'),
        # 通用容器
        ('article', None),
        ('div', 'step'),
        ('div', 'card'),
    ]
    
    # 需要过滤的无效slide特征
    SKIP_PATTERNS = [
        r'^\s*$',  # 空白内容
        r'^<!--.*-->$',  # 纯注释
    ]
    
    def __init__(self, html_path: str, output_path: str = None, callback=None):
        self.html_path = Path(html_path).resolve()
        if not self.html_path.exists():
            raise FileNotFoundError(f"文件不存在: {html_path}")
        
        self.output_path = output_path or str(self.html_path.with_suffix('.pptx'))
        self.temp_dir = tempfile.mkdtemp(prefix="csndc_ppt_")
        self.callback = callback or (lambda msg: None)
        
        from bs4 import BeautifulSoup
        from pptx import Presentation
        from pptx.util import Inches
        
        # 读取并解析HTML
        self.html_content = self.html_path.read_text(encoding='utf-8')
        self.soup = BeautifulSoup(self.html_content, 'html.parser')
        
        # 提取所有样式（内联+外部引用的CSS变量）
        self.css = self._extract_styles()
        
        # 初始化PPT（16:9）
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
    
    def log(self, msg: str):
        self.callback(msg)
    
    def _extract_styles(self) -> str:
        """提取所有样式内容"""
        styles = []
        
        # 内联style标签
        for style in self.soup.find_all('style'):
            styles.append(style.get_text())
        
        # 提取CSS变量定义（:root）
        root_vars = re.findall(r':root\s*\{[^}]+\}', self.html_content)
        styles.extend(root_vars)
        
        return '\n'.join(styles)
    
    def _is_valid_slide(self, elem) -> bool:
        """检查元素是否是有效的slide"""
        if elem is None:
            return False
        
        # 获取文本内容
        text = elem.get_text(strip=True)
        
        # 跳过空内容
        if not text and not elem.find(['img', 'svg', 'canvas', 'video']):
            return False
        
        # 跳过隐藏元素
        style = elem.get('style', '')
        if 'display:none' in style.replace(' ', '') or 'visibility:hidden' in style.replace(' ', ''):
            return False
        
        # 跳过过小的容器
        if len(text) < 5 and not elem.find(['img', 'svg', 'canvas']):
            return False
        
        return True
    
    def _find_slides(self) -> Tuple[List, str]:
        """智能检测HTML中的slide元素
        
        Returns:
            (slides列表, 检测到的选择器描述)
        """
        for tag, cls in self.SLIDE_SELECTORS:
            if cls:
                candidates = self.soup.find_all(tag, class_=lambda c: c and cls in c.split())
            else:
                candidates = self.soup.find_all(tag)
            
            # 过滤无效slide
            slides = [s for s in candidates if self._is_valid_slide(s)]
            
            if len(slides) >= 1:
                selector_desc = f"<{tag} class='{cls}'" if cls else f"<{tag}>"
                self.log(f"检测到 {len(slides)} 个slide ({selector_desc})")
                return slides, selector_desc
        
        # 未找到slide结构，整页转换
        body = self.soup.find('body')
        if body:
            self.log("未检测到slide结构，整页转为1张")
            return [body], "<body>"
        
        return [], ""
    
    def _get_title(self, elem) -> str:
        """智能提取slide标题"""
        # 按优先级查找标题
        for selector in ['h1', 'h2', 'h3', '.title', '.header', '[class*="title"]']:
            if selector.startswith('.'):
                t = elem.find(class_=selector[1:])
            elif selector.startswith('['):
                # 属性选择器
                t = elem.find(attrs={'class': re.compile(r'title', re.I)})
            else:
                t = elem.find(selector)
            
            if t:
                text = t.get_text(strip=True)[:40]
                if text:
                    return text
        
        # 取第一段文字作为标题
        first_text = elem.get_text(strip=True)[:40]
        return first_text if first_text else ""
    
    def _build_slide_html(self, slide_elem, is_full_page: bool = False) -> str:
        """构建渲染用HTML
        
        增强功能：
        - 保持原始样式
        - 隐藏滚动条
        - 固定viewport尺寸
        - 处理特殊布局
        """
        # 通用样式注入
        inject_css = f"""
/* CSNDC PPT Converter - 注入样式 */
* {{
    scrollbar-width: none !important;
    -ms-overflow-style: none !important;
}}
*::-webkit-scrollbar {{
    display: none !important;
}}
html, body {{
    margin: 0 !important;
    padding: 0 !important;
    overflow: hidden !important;
    background: white;
}}
/* 动画冻结 */
*, *::before, *::after {{
    animation-play-state: paused !important;
    transition: none !important;
}}
"""
        
        if is_full_page:
            return f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8">
<style>
{self.css}
{inject_css}
</style></head>
<body>{slide_elem}</body></html>"""
        
        # Slide模式：精确控制尺寸
        slide_css = f"""
{inject_css}
html, body {{
    width: {self.VIEWPORT_W}px !important;
    height: {self.VIEWPORT_H}px !important;
}}
/* 通用slide容器适配 */
.slide, .page, section, article, .step, .swiper-slide, .carousel-item {{
    margin: 0 !important;
    box-shadow: none !important;
    border-radius: 0 !important;
    width: {self.VIEWPORT_W}px !important;
    height: {self.VIEWPORT_H}px !important;
    min-height: {self.VIEWPORT_H}px !important;
    max-height: {self.VIEWPORT_H}px !important;
    box-sizing: border-box !important;
    position: absolute !important;
    top: 0 !important;
    left: 0 !important;
    overflow: hidden !important;
}}
"""
        
        return f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8">
<style>
{self.css}
{slide_css}
</style></head>
<body>{slide_elem}</body></html>"""
    
    async def _render_all(self, slides: List) -> List[str]:
        """批量渲染slides（复用浏览器实例）"""
        from playwright.async_api import async_playwright
        
        img_paths = []
        is_full_page = len(slides) == 1 and slides[0].name == 'body'
        
        self.log("启动渲染引擎...")
        
        async with async_playwright() as p:
            # 浏览器优先级: Edge > Chrome > Chromium
            browser = None
            browser_name = ""
            last_error = None
            
            for channel in ['msedge', 'chrome', None]:
                try:
                    channel_name = channel or 'chromium'
                    self.log(f"尝试连接: {channel_name}")
                    
                    if channel:
                        browser = await p.chromium.launch(
                            headless=True, 
                            channel=channel,
                            args=['--disable-gpu', '--no-sandbox', '--disable-dev-shm-usage']
                        )
                        browser_name = "Edge" if channel == 'msedge' else "Chrome"
                    else:
                        browser = await p.chromium.launch(
                            headless=True,
                            args=['--disable-gpu', '--no-sandbox', '--disable-dev-shm-usage']
                        )
                        browser_name = "Chromium"
                    
                    self.log(f"✓ 使用渲染引擎: {browser_name}")
                    break
                except Exception as e:
                    last_error = str(e)
                    self.log(f"  {channel_name} 不可用")
                    continue
            
            if not browser:
                raise RuntimeError(f"未找到可用浏览器\n{last_error}")
            
            context = await browser.new_context(
                viewport={'width': self.VIEWPORT_W, 'height': self.VIEWPORT_H},
                device_scale_factor=self.SCALE
            )
            page = await context.new_page()
            
            total = len(slides)
            for i, elem in enumerate(slides):
                num = i + 1
                title = self._get_title(elem) or f"Slide {num}"
                self.log(f"[{num}/{total}] {title}")
                
                # 写入临时HTML
                html_file = Path(self.temp_dir) / f"slide_{num}.html"
                html_content = self._build_slide_html(elem, is_full_page)
                html_file.write_text(html_content, encoding='utf-8')
                
                # 渲染
                img_path = Path(self.temp_dir) / f"slide_{num}.png"
                await page.goto(f'file:///{html_file}')
                
                # 等待加载完成
                await page.wait_for_load_state('networkidle')
                await page.wait_for_timeout(200)  # 额外等待渲染
                
                # 截图
                await page.screenshot(path=str(img_path), full_page=is_full_page)
                
                if img_path.exists():
                    img_paths.append(str(img_path))
            
            await browser.close()
            self.log(f"渲染完成: {len(img_paths)}张图片")
        
        return img_paths
    
    def _add_slide(self, img_path: str) -> bool:
        """添加图片到PPT"""
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        
        try:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
            slide.shapes.add_picture(
                img_path, Inches(0), Inches(0),
                self.prs.slide_width, self.prs.slide_height
            )
            return True
        except Exception as e:
            self.log(f"添加失败: {e}")
            return False
    
    def convert(self) -> str:
        """执行转换"""
        self.log(f"输入: {self.html_path.name}")
        self.log(f"输出分辨率: 4K ({self.VIEWPORT_W * self.SCALE}×{self.VIEWPORT_H * self.SCALE})")
        
        # 检测slides
        slides, _ = self._find_slides()
        if not slides:
            raise ValueError("未找到可转换的内容")
        
        # 渲染
        self.log("开始渲染...")
        img_paths = asyncio.run(self._render_all(slides))
        
        # 生成PPT
        self.log("生成PPT...")
        success = 0
        for img_path in img_paths:
            if self._add_slide(img_path):
                success += 1
        
        self.prs.save(self.output_path)
        
        # 清理临时文件
        shutil.rmtree(self.temp_dir, ignore_errors=True)
        
        # 结果
        size_mb = Path(self.output_path).stat().st_size / (1024 * 1024)
        self.log(f"")
        self.log(f"✓ 转换完成: {success}/{len(slides)} slides")
        self.log(f"  输出: {Path(self.output_path).name} ({size_mb:.1f}MB)")
        
        return self.output_path


# ============== GUI 应用 ==============

# 联想品牌色系 - 现代科技感
class Colors:
    RED = "#E2231A"           # 联想红 - 主色
    RED_DARK = "#B81C15"      # 深红 - 悬停态
    RED_LIGHT = "#FF4D4D"     # 浅红 - 高亮
    
    BG_DARK = "#0D0D0D"       # 深黑背景
    BG_CARD = "#1A1A1A"       # 卡片背景
    BG_HOVER = "#252525"      # 悬停背景
    
    TEXT_PRIMARY = "#FFFFFF"  # 主文字
    TEXT_SECONDARY = "#888888" # 次要文字
    TEXT_MUTED = "#555555"    # 暗文字
    
    BORDER = "#333333"        # 边框
    ACCENT = "#E2231A"        # 强调色


class ModernButton(tk.Canvas):
    """现代化按钮组件"""
    
    def __init__(self, parent, text, command=None, width=200, height=45, **kwargs):
        super().__init__(parent, width=width, height=height, 
                        bg=Colors.BG_DARK, highlightthickness=0, **kwargs)
        
        self.command = command
        self.text = text
        self.width = width
        self.height = height
        self.is_hover = False
        
        self._draw()
        self.bind("<Enter>", self._on_enter)
        self.bind("<Leave>", self._on_leave)
        self.bind("<Button-1>", self._on_click)
    
    def _draw(self):
        self.delete("all")
        
        # 背景色
        bg = Colors.RED_DARK if self.is_hover else Colors.RED
        
        # 绘制圆角矩形
        r = 8  # 圆角半径
        self.create_arc(0, 0, r*2, r*2, start=90, extent=90, fill=bg, outline=bg)
        self.create_arc(self.width-r*2, 0, self.width, r*2, start=0, extent=90, fill=bg, outline=bg)
        self.create_arc(0, self.height-r*2, r*2, self.height, start=180, extent=90, fill=bg, outline=bg)
        self.create_arc(self.width-r*2, self.height-r*2, self.width, self.height, start=270, extent=90, fill=bg, outline=bg)
        self.create_rectangle(r, 0, self.width-r, self.height, fill=bg, outline=bg)
        self.create_rectangle(0, r, self.width, self.height-r, fill=bg, outline=bg)
        
        # 文字
        self.create_text(
            self.width/2, self.height/2,
            text=self.text,
            font=("Segoe UI Semibold", 11),
            fill=Colors.TEXT_PRIMARY
        )
    
    def _on_enter(self, e):
        self.is_hover = True
        self._draw()
        self.config(cursor="hand2")
    
    def _on_leave(self, e):
        self.is_hover = False
        self._draw()
    
    def _on_click(self, e):
        if self.command:
            self.command()


class DropZone(tk.Canvas):
    """拖放区域组件 - 虚线边框+动态效果"""
    
    def __init__(self, parent, on_click=None, **kwargs):
        super().__init__(parent, highlightthickness=0, bg=Colors.BG_DARK, **kwargs)
        
        self.on_click = on_click
        self.is_hover = False
        self.is_processing = False
        self.status_text = "点击选择文件"
        self.sub_text = "支持 .html / .htm 格式"
        
        self.bind("<Configure>", self._on_resize)
        self.bind("<Enter>", self._on_enter)
        self.bind("<Leave>", self._on_leave)
        self.bind("<Button-1>", self._on_click)
    
    def _draw(self):
        self.delete("all")
        w, h = self.winfo_width(), self.winfo_height()
        if w < 10 or h < 10:
            return
        
        # 背景
        bg = Colors.BG_HOVER if self.is_hover else Colors.BG_CARD
        self._draw_rounded_rect(4, 4, w-4, h-4, 16, fill=bg, outline="")
        
        # 虚线边框
        border_color = Colors.RED if self.is_hover else Colors.BORDER
        self._draw_dashed_border(12, 12, w-12, h-12, 12, border_color)
        
        # 图标
        icon_y = h/2 - 35
        if self.is_processing:
            # 转换中图标
            self.create_text(w/2, icon_y, text="⚙", font=("Segoe UI Symbol", 42), fill=Colors.RED)
        else:
            # 文件图标
            self._draw_file_icon(w/2, icon_y)
        
        # 主文字
        text_color = Colors.RED if self.is_processing else Colors.TEXT_PRIMARY
        self.create_text(
            w/2, h/2 + 20,
            text=self.status_text,
            font=("微软雅黑", 13, "bold"),
            fill=text_color
        )
        
        # 副文字
        self.create_text(
            w/2, h/2 + 50,
            text=self.sub_text,
            font=("Segoe UI", 9),
            fill=Colors.TEXT_MUTED
        )
    
    def _draw_rounded_rect(self, x1, y1, x2, y2, r, **kwargs):
        """ 绘制圆角矩形"""
        points = [
            x1+r, y1, x2-r, y1, x2, y1, x2, y1+r,
            x2, y2-r, x2, y2, x2-r, y2, x1+r, y2,
            x1, y2, x1, y2-r, x1, y1+r, x1, y1
        ]
        self.create_polygon(points, smooth=True, **kwargs)
    
    def _draw_dashed_border(self, x1, y1, x2, y2, r, color):
        """绘制虚线边框"""
        dash = (8, 4)
        # 上边
        self.create_line(x1+r, y1, x2-r, y1, fill=color, dash=dash, width=2)
        # 下边
        self.create_line(x1+r, y2, x2-r, y2, fill=color, dash=dash, width=2)
        # 左边
        self.create_line(x1, y1+r, x1, y2-r, fill=color, dash=dash, width=2)
        # 右边
        self.create_line(x2, y1+r, x2, y2-r, fill=color, dash=dash, width=2)
        # 圆角
        self.create_arc(x1, y1, x1+r*2, y1+r*2, start=90, extent=90, style="arc", outline=color, width=2)
        self.create_arc(x2-r*2, y1, x2, y1+r*2, start=0, extent=90, style="arc", outline=color, width=2)
        self.create_arc(x1, y2-r*2, x1+r*2, y2, start=180, extent=90, style="arc", outline=color, width=2)
        self.create_arc(x2-r*2, y2-r*2, x2, y2, start=270, extent=90, style="arc", outline=color, width=2)
    
    def _draw_file_icon(self, cx, cy):
        """绘制文件图标"""
        # 文件外框
        size = 28
        fold = 10
        color = Colors.RED if self.is_hover else Colors.TEXT_SECONDARY
        
        points = [
            cx-size, cy-size,
            cx+size-fold, cy-size,
            cx+size, cy-size+fold,
            cx+size, cy+size,
            cx-size, cy+size
        ]
        self.create_polygon(points, fill="", outline=color, width=2)
        
        # 折角
        self.create_line(cx+size-fold, cy-size, cx+size-fold, cy-size+fold, fill=color, width=2)
        self.create_line(cx+size-fold, cy-size+fold, cx+size, cy-size+fold, fill=color, width=2)
        
        # HTML文字
        self.create_text(cx, cy+5, text="HTML", font=("Consolas", 10, "bold"), fill=color)
    
    def _on_resize(self, e):
        self._draw()
    
    def _on_enter(self, e):
        if not self.is_processing:
            self.is_hover = True
            self._draw()
            self.config(cursor="hand2")
    
    def _on_leave(self, e):
        self.is_hover = False
        self._draw()
    
    def _on_click(self, e):
        if self.on_click and not self.is_processing:
            self.on_click()
    
    def set_processing(self, processing: bool, text: str = None):
        self.is_processing = processing
        if text:
            self.status_text = text
        else:
            self.status_text = "正在转换..." if processing else "点击选择文件"
        self.sub_text = "请稍候" if processing else "支持 .html / .htm 格式"
        self._draw()


class LogPanel(tk.Frame):
    """日志面板 - 终端风格"""
    
    def __init__(self, parent, **kwargs):
        super().__init__(parent, bg=Colors.BG_DARK, **kwargs)
        
        # 标题栏
        header = tk.Frame(self, bg=Colors.BG_CARD, height=32)
        header.pack(fill=tk.X)
        header.pack_propagate(False)
        
        # 终端指示点
        dots_frame = tk.Frame(header, bg=Colors.BG_CARD)
        dots_frame.pack(side=tk.LEFT, padx=12, pady=10)
        
        for color in ["#FF5F56", "#FFBD2E", "#27CA40"]:
            dot = tk.Canvas(dots_frame, width=12, height=12, bg=Colors.BG_CARD, highlightthickness=0)
            dot.create_oval(0, 0, 12, 12, fill=color, outline="")
            dot.pack(side=tk.LEFT, padx=2)
        
        tk.Label(
            header, text="Console Output", 
            font=("Consolas", 9), 
            bg=Colors.BG_CARD, 
            fg=Colors.TEXT_MUTED
        ).pack(side=tk.LEFT, padx=10)
        
        # 日志文本区 - 增大高度和字体
        self.text = tk.Text(
            self,
            font=("Consolas", 11),  # 增大字体
            bg=Colors.BG_CARD,
            fg=Colors.TEXT_SECONDARY,
            insertbackground=Colors.RED,
            selectbackground=Colors.RED,
            relief="flat",
            padx=15,
            pady=12,
            height=12,  # 增大高度，可显示更多行
            spacing1=4,  # 行间距
            state="disabled",
            cursor="arrow"
        )
        self.text.pack(fill=tk.BOTH, expand=True, padx=1, pady=(0, 1))
        
        # 配置标签样式
        self.text.tag_configure("success", foreground="#27CA40")      # 绿色 - 成功
        self.text.tag_configure("error", foreground="#FF5F56")        # 红色 - 错误
        self.text.tag_configure("info", foreground=Colors.TEXT_SECONDARY)  # 灰色 - 普通信息
        self.text.tag_configure("progress", foreground="#5AC8FA")     # 蓝色 - 进度
        self.text.tag_configure("highlight", foreground="#FFBD2E")    # 黄色 - 高亮提示
    
    def log(self, msg: str, tag: str = "info"):
        self.text.config(state="normal")
        
        # 根据内容自动选择标签
        if "✓" in msg or "完成" in msg or "成功" in msg:
            tag = "success"
        elif "错误" in msg or "失败" in msg or "Error" in msg:
            tag = "error"
        elif msg.startswith("[") and "/" in msg:  # [1/10] 这种进度格式
            tag = "progress"
        elif "引擎" in msg or "输出" in msg or "输入" in msg:
            tag = "highlight"
        
        self.text.insert(tk.END, f"> {msg}\n", tag)
        self.text.see(tk.END)
        self.text.config(state="disabled")
    
    def clear(self):
        self.text.config(state="normal")
        self.text.delete(1.0, tk.END)
        self.text.config(state="disabled")


class ProgressBar(tk.Canvas):
    """自定义进度条 - 动画效果"""
    
    def __init__(self, parent, **kwargs):
        super().__init__(parent, height=4, bg=Colors.BG_DARK, highlightthickness=0, **kwargs)
        
        self.progress = 0
        self.animating = False
        self.anim_pos = 0
        
        self.bind("<Configure>", self._draw)
    
    def _draw(self, event=None):
        self.delete("all")
        w = self.winfo_width()
        h = self.winfo_height()
        
        # 背景轨道
        self.create_rectangle(0, 0, w, h, fill=Colors.BORDER, outline="")
        
        if self.animating:
            # 动画模式 - 滑动块
            block_w = w * 0.3
            x = (self.anim_pos % 100) / 100 * (w + block_w) - block_w
            self.create_rectangle(x, 0, x + block_w, h, fill=Colors.RED, outline="")
    
    def start(self):
        self.animating = True
        self._animate()
    
    def stop(self):
        self.animating = False
        self._draw()
    
    def _animate(self):
        if self.animating:
            self.anim_pos = (self.anim_pos + 2) % 100
            self._draw()
            self.after(20, self._animate)


class App(tk.Tk):
    """CSNDC PPT Converter - 现代化GUI"""
    
    def __init__(self):
        super().__init__()
        
        self.title("CSNDC PPT Converter")
        self.geometry("600x620")  # 增大窗口高度
        self.resizable(False, False)
        self.configure(bg=Colors.BG_DARK)
        
        # 居中显示
        self.update_idletasks()
        x = (self.winfo_screenwidth() - 600) // 2
        y = (self.winfo_screenheight() - 620) // 2
        self.geometry(f"+{x}+{y}")
        
        self._setup_ui()
    
    def _setup_ui(self):
        """构建现代化界面"""
        
        # ===== 顶部标题区 =====
        header = tk.Frame(self, bg=Colors.BG_DARK, height=80)
        header.pack(fill=tk.X, padx=30, pady=(25, 0))
        header.pack_propagate(False)
        
        # Logo区域
        logo_frame = tk.Frame(header, bg=Colors.BG_DARK)
        logo_frame.pack(side=tk.LEFT)
        
        # 红色标识块
        logo_canvas = tk.Canvas(logo_frame, width=8, height=50, bg=Colors.BG_DARK, highlightthickness=0)
        logo_canvas.pack(side=tk.LEFT, padx=(0, 15))
        logo_canvas.create_rectangle(0, 0, 8, 50, fill=Colors.RED, outline="")
        
        # 标题文字
        title_frame = tk.Frame(logo_frame, bg=Colors.BG_DARK)
        title_frame.pack(side=tk.LEFT)
        
        tk.Label(
            title_frame,
            text="CSNDC PPT Converter",
            font=("Segoe UI", 22, "bold"),
            bg=Colors.BG_DARK,
            fg=Colors.TEXT_PRIMARY
        ).pack(anchor="w")
        
        tk.Label(
            title_frame,
            text="HTML → PowerPoint · 4K Ultra HD",
            font=("Segoe UI", 10),
            bg=Colors.BG_DARK,
            fg=Colors.TEXT_MUTED
        ).pack(anchor="w", pady=(2, 0))
        
        # 右侧By Dolly
        dolly_frame = tk.Frame(header, bg=Colors.BG_DARK)
        dolly_frame.pack(side=tk.RIGHT, pady=10)
        
        tk.Label(
            dolly_frame,
            text="By",
            font=("Segoe UI", 9),
            bg=Colors.BG_DARK,
            fg=Colors.TEXT_MUTED
        ).pack()
        
        tk.Label(
            dolly_frame,
            text="Dolly",
            font=("Segoe UI Semibold", 14),
            bg=Colors.BG_DARK,
            fg=Colors.RED
        ).pack()
        
        # ===== 拖放区 =====
        self.drop_zone = DropZone(self, on_click=self._select_file)
        self.drop_zone.pack(fill=tk.BOTH, expand=True, padx=30, pady=20)
        
        # ===== 日志区 =====
        self.log_panel = LogPanel(self)
        self.log_panel.pack(fill=tk.X, padx=30, pady=(0, 15))
        
        # ===== 进度条 =====
        self.progress = ProgressBar(self)
        self.progress.pack(fill=tk.X, padx=30)
        
        # ===== 底部信息 =====
        footer = tk.Frame(self, bg=Colors.BG_DARK, height=40)
        footer.pack(fill=tk.X, padx=30, pady=(15, 20))
        
        # 左侧状态
        tk.Label(
            footer,
            text="● Ready",
            font=("Segoe UI", 9),
            bg=Colors.BG_DARK,
            fg="#27CA40"
        ).pack(side=tk.LEFT)
        
        # 右侧版本
        tk.Label(
            footer,
            text="v1.0 · Lenovo CSNDC",
            font=("Segoe UI", 9),
            bg=Colors.BG_DARK,
            fg=Colors.TEXT_MUTED
        ).pack(side=tk.RIGHT)
    
    def _select_file(self):
        """选择文件"""
        file_path = filedialog.askopenfilename(
            title="选择HTML文件",
            filetypes=[("HTML文件", "*.html *.htm"), ("所有文件", "*.*")]
        )
        if file_path:
            self._process_file(file_path)
    
    def _log(self, msg: str):
        """输出日志"""
        self.log_panel.log(msg)
        self.update_idletasks()
    
    def _process_file(self, file_path: str):
        """处理文件转换"""
        if not file_path.lower().endswith(('.html', '.htm')):
            messagebox.showerror("错误", "请选择HTML文件")
            return
        
        if not os.path.exists(file_path):
            messagebox.showerror("错误", "文件不存在")
            return
        
        self.log_panel.clear()
        self.progress.start()
        self.drop_zone.set_processing(True)
        
        def run():
            try:
                converter = HTMLToPPT(
                    file_path, 
                    callback=lambda m: self.after(0, lambda: self._log(m))
                )
                output = converter.convert()
                self.after(0, lambda: self._on_complete(output))
            except Exception as e:
                self.after(0, lambda: self._on_error(str(e)))
        
        threading.Thread(target=run, daemon=True).start()
    
    def _on_complete(self, output: str):
        """转换完成"""
        self.progress.stop()
        self.drop_zone.set_processing(False)
        
        result = messagebox.askyesno(
            "转换完成", 
            f"已生成: {Path(output).name}\n\n是否打开文件所在目录？"
        )
        if result:
            os.startfile(str(Path(output).parent))
    
    def _on_error(self, msg: str):
        """转换失败"""
        self.progress.stop()
        self.drop_zone.set_processing(False)
        self._log(f"错误: {msg}")
        messagebox.showerror("转换失败", msg)


def main():
    missing = check_dependencies()
    if missing:
        root = tk.Tk()
        root.withdraw()
        messagebox.showerror(
            "缺少依赖",
            f"请安装以下依赖:\npip install {' '.join(missing)}"
            + ("\n\n还需运行:\npython -m playwright install" if "playwright" in missing else "")
        )
        return 1
    
    app = App()
    app.mainloop()
    return 0


if __name__ == "__main__":
    sys.exit(main())

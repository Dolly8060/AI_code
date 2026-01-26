"""HTML to PPT Converter - 4K高清通用版

通用HTML转PPT工具，支持：
1. 多slide HTML（自动检测 .slide / .page / section 等元素）
2. 单页面HTML（整页转为一张slide）
3. 4K高清输出（4000×2250）

使用: python ppt.py input.html [-o output.pptx]
"""

import os
import sys
import asyncio
import tempfile
import shutil
import argparse
from pathlib import Path
from typing import List, Optional

# 依赖检查
def check_deps():
    missing = []
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        missing.append("beautifulsoup4")
    try:
        from pptx import Presentation
    except ImportError:
        missing.append("python-pptx")
    try:
        from playwright.async_api import async_playwright
    except ImportError:
        missing.append("playwright")
    try:
        from PIL import Image
    except ImportError:
        missing.append("Pillow")
    
    if missing:
        print(f"缺少依赖: {', '.join(missing)}")
        print(f"安装命令: pip install {' '.join(missing)}")
        if "playwright" in missing:
            print("还需运行: python -m playwright install chromium")
        sys.exit(1)

check_deps()

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from playwright.async_api import async_playwright
from PIL import Image


class HTMLToPPT:
    """通用HTML转PPT转换器（4K高清）"""
    
    # 4K配置：viewport 1000x563 + 4x缩放 = 4000x2252
    VIEWPORT_W = 1000
    VIEWPORT_H = 563
    SCALE = 4
    
    # PPT尺寸（16:9）
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)
    
    # 常见slide选择器（优先级从高到低）
    SLIDE_SELECTORS = [
        ('div', 'slide'),
        ('section', 'slide'),
        ('div', 'page'),
        ('section', None),  # 纯section标签
    ]
    
    def __init__(self, html_path: str, output_path: str = None):
        self.html_path = Path(html_path).resolve()
        if not self.html_path.exists():
            raise FileNotFoundError(f"文件不存在: {html_path}")
        
        self.output_path = output_path or str(self.html_path.with_suffix('.pptx'))
        self.temp_dir = tempfile.mkdtemp(prefix="html2ppt_")
        
        # 读取HTML
        self.html_content = self.html_path.read_text(encoding='utf-8')
        self.soup = BeautifulSoup(self.html_content, 'html.parser')
        
        # 提取样式
        self.css = '\n'.join(s.get_text() for s in self.soup.find_all('style'))
        
        # 初始化PPT
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT
    
    def _find_slides(self) -> List:
        """自动检测HTML中的slide元素"""
        for tag, cls in self.SLIDE_SELECTORS:
            if cls:
                slides = self.soup.find_all(tag, class_=cls)
            else:
                slides = self.soup.find_all(tag)
            if slides:
                print(f"检测到 {len(slides)} 个slide (<{tag} class='{cls or ''}'>)")
                return slides
        
        # 未找到slide结构，将整个body作为单页
        body = self.soup.find('body')
        if body:
            print("未检测到slide结构，将整页转为1张slide")
            return [body]
        return []
    
    def _get_title(self, elem) -> str:
        """提取slide标题"""
        for tag in ['h1', 'h2', 'h3', '.title', '.header']:
            title_elem = elem.select_one(tag) if tag.startswith('.') else elem.find(tag)
            if title_elem:
                text = title_elem.get_text(strip=True)[:40]
                if text:
                    return text
        return ""
    
    def _build_slide_html(self, slide_elem, is_full_page: bool = False) -> str:
        """构建单个slide的完整HTML"""
        if is_full_page:
            # 整页模式：直接使用原始HTML
            return f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8">
<style>
{self.css}
* {{ scrollbar-width: none !important; -ms-overflow-style: none !important; }}
*::-webkit-scrollbar {{ display: none !important; }}
html, body {{ margin: 0; padding: 0; overflow: hidden; }}
</style></head>
<body>{str(slide_elem)}</body></html>"""
        
        # Slide模式：精确控制尺寸
        return f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8">
<style>
{self.css}
* {{ scrollbar-width: none !important; -ms-overflow-style: none !important; }}
*::-webkit-scrollbar {{ display: none !important; }}
html, body {{
    margin: 0 !important; padding: 0 !important;
    width: {self.VIEWPORT_W}px !important;
    height: {self.VIEWPORT_H}px !important;
    overflow: hidden !important;
    background: white;
}}
.slide, .page, section {{
    margin: 0 !important; padding: 20px !important;
    box-shadow: none !important; border-radius: 0 !important;
    width: {self.VIEWPORT_W}px !important;
    height: {self.VIEWPORT_H}px !important;
    box-sizing: border-box !important;
    position: absolute !important; top: 0 !important; left: 0 !important;
}}
</style></head>
<body>{str(slide_elem)}</body></html>"""
    
    async def _render_all(self, slides: List) -> List[str]:
        """批量渲染所有slides（复用browser实例）"""
        img_paths = []
        is_full_page = len(slides) == 1 and slides[0].name == 'body'
        
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            context = await browser.new_context(
                viewport={'width': self.VIEWPORT_W, 'height': self.VIEWPORT_H},
                device_scale_factor=self.SCALE
            )
            page = await context.new_page()
            
            for i, slide_elem in enumerate(slides):
                num = i + 1
                title = self._get_title(slide_elem) or f"Slide {num}"
                print(f"  [{num}/{len(slides)}] {title}", end="", flush=True)
                
                # 写入临时HTML
                html_file = Path(self.temp_dir) / f"slide_{num}.html"
                html_file.write_text(self._build_slide_html(slide_elem, is_full_page), encoding='utf-8')
                
                # 渲染
                img_path = Path(self.temp_dir) / f"slide_{num}.png"
                await page.goto(f'file:///{html_file}')
                await page.wait_for_load_state('networkidle')
                await page.screenshot(path=str(img_path), full_page=is_full_page)
                
                if img_path.exists():
                    with Image.open(img_path) as img:
                        w, h = img.size
                    size_kb = img_path.stat().st_size / 1024
                    print(f" -> {w}x{h} ({size_kb:.0f}KB) ✓")
                    img_paths.append(str(img_path))
                else:
                    print(" ✗")
            
            await browser.close()
        
        return img_paths
    
    def _add_slide(self, img_path: str) -> bool:
        """添加图片到PPT"""
        try:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
            slide.shapes.add_picture(
                img_path, Inches(0), Inches(0),
                self.SLIDE_WIDTH, self.SLIDE_HEIGHT
            )
            return True
        except Exception as e:
            print(f"添加slide失败: {e}")
            return False
    
    def convert(self) -> str:
        """执行转换"""
        print(f"\n{'='*50}")
        print(f"HTML转PPT (4K: {self.VIEWPORT_W * self.SCALE}x{self.VIEWPORT_H * self.SCALE})")
        print(f"输入: {self.html_path.name}")
        print(f"{'='*50}")
        
        # 检测slides
        slides = self._find_slides()
        if not slides:
            raise ValueError("未找到可转换的内容")
        
        # 渲染
        print("\n渲染中...")
        img_paths = asyncio.run(self._render_all(slides))
        
        # 生成PPT
        print("\n生成PPT...")
        success = 0
        for img_path in img_paths:
            if self._add_slide(img_path):
                success += 1
        
        self.prs.save(self.output_path)
        
        # 清理
        shutil.rmtree(self.temp_dir, ignore_errors=True)
        
        # 结果
        ppt_size = Path(self.output_path).stat().st_size / (1024 * 1024)
        print(f"\n{'='*50}")
        print(f"✓ 完成: {success}/{len(slides)} slides")
        print(f"  输出: {self.output_path}")
        print(f"  大小: {ppt_size:.1f} MB")
        print(f"{'='*50}")
        
        return self.output_path


def main():
    parser = argparse.ArgumentParser(
        description='HTML转PPT（4K高清）',
        epilog='示例: python ppt.py slide.html -o output.pptx'
    )
    parser.add_argument('html_file', help='HTML文件路径')
    parser.add_argument('-o', '--output', help='输出PPT路径（默认同名.pptx）')
    
    args = parser.parse_args()
    
    try:
        converter = HTMLToPPT(args.html_file, args.output)
        converter.convert()
        return 0
    except Exception as e:
        print(f"\n错误: {e}")
        return 1


if __name__ == '__main__':
    sys.exit(main())

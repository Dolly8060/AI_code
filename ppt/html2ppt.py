"""HTML to PPT Converter - 专门针对 ai_zone.html 的高保真转换工具

特点：
- 1:1 完美复刻 HTML 视觉效果
- 支持 4K/8K/16K 等多种分辨率选项
- 保留所有 CSS 样式、SVG 图形、复杂布局
- 16:9 标准 PPT 比例输出
"""

import os
import asyncio
import tempfile
import shutil
from pathlib import Path

# 依赖检查
try:
    from bs4 import BeautifulSoup
except ImportError:
    print("请安装: pip install beautifulsoup4")
    exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("请安装: pip install python-pptx")
    exit(1)

try:
    from playwright.async_api import async_playwright
except ImportError:
    print("请安装: pip install playwright && python -m playwright install chromium")
    exit(1)

try:
    from PIL import Image
except ImportError:
    print("请安装: pip install Pillow")
    exit(1)


class AIZoneHTMLToPPT:
    """专门针对 ai_zone.html 的高保真 PPT 转换器"""
    
    # HTML 原始尺寸 (16:9 比例)
    HTML_WIDTH = 1000
    HTML_HEIGHT = 562.5
    
    # 分辨率预设: (viewport_width, viewport_height, dpi_scale, description)
    # 策略: DPI 保持在 4x 以下确保字体清晰，通过增大 viewport 达到更高分辨率
    RESOLUTIONS = {
        '1080p': (1000, 563, 2, '2000x1126'),      # 2x DPI
        '2k':    (1000, 563, 3, '3000x1689'),      # 3x DPI  
        '4k':    (1000, 563, 4, '4000x2252'),      # 4x DPI [推荐，最清晰]
        '5k':    (1280, 720, 4, '5120x2880'),      # viewport放大 + 4x DPI
        '8k':    (1920, 1080, 4, '7680x4320'),     # viewport放大 + 4x DPI
        '16k':   (3840, 2160, 4, '15360x8640'),    # viewport放大 + 4x DPI [超高清]
    }
    
    # 颜色主题 (从 CSS 变量提取)
    COLORS = {
        'lenovo_green': '#58816F',
        'lenovo_dark': '#333333',
        'lenovo_light': '#F4F8F6',
        'accent_green': '#7CAD95',
        'code_bg': '#1E1E1E',
        'text_grey': '#555555',
    }
    
    def __init__(self, html_path: str = None, output_path: str = None, resolution: str = '4k'):
        """初始化转换器
        
        Args:
            html_path: HTML 文件路径，默认为 ai_zone_v2.html
            output_path: 输出 PPT 路径，默认为 ai_zone_v2.pptx
            resolution: 分辨率选项 (1080p/2k/4k/5k/8k/16k)，默认 4k
        """
        # 默认路径
        workspace = Path(__file__).parent
        self.html_path = Path(html_path) if html_path else workspace / "ai_zone_v2.html"
        self.output_path = Path(output_path) if output_path else workspace / "ai_zone_v2.pptx"
        
        # 验证文件存在
        if not self.html_path.exists():
            raise FileNotFoundError(f"HTML 文件不存在: {self.html_path}")
        
        # 解析分辨率配置
        self.resolution = resolution.lower()
        if self.resolution in self.RESOLUTIONS:
            self.viewport_w, self.viewport_h, self.scale_factor, self.res_desc = self.RESOLUTIONS[self.resolution]
        else:
            # 支持自定义 DPI (如 "6x" 表示 6 倍 DPI)
            if self.resolution.endswith('x') and self.resolution[:-1].isdigit():
                custom_scale = int(self.resolution[:-1])
                self.viewport_w, self.viewport_h = int(self.HTML_WIDTH), int(self.HTML_HEIGHT)
                self.scale_factor = min(custom_scale, 16)  # 最大 16x
                self.res_desc = f'{self.viewport_w * self.scale_factor}x{self.viewport_h * self.scale_factor}'
            else:
                print(f"⚠️ 未知分辨率 '{resolution}'，使用默认 4k")
                self.viewport_w, self.viewport_h, self.scale_factor, self.res_desc = self.RESOLUTIONS['4k']
                self.resolution = '4k'
        
        # 临时目录
        self.temp_dir = tempfile.mkdtemp(prefix="html2ppt_")
        
        # 计算渲染尺寸
        self.render_width = int(self.viewport_w * self.scale_factor)
        self.render_height = int(self.viewport_h * self.scale_factor)
        
        # CSS 缩放比例（viewport 放大时需要缩放原始内容）
        self.css_scale = self.viewport_w / self.HTML_WIDTH
        
        # 初始化 PPT (16:9 宽屏)
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 标准 16:9
        self.prs.slide_height = Inches(7.5)
        
        # 读取并解析 HTML
        with open(self.html_path, 'r', encoding='utf-8') as f:
            self.html_content = f.read()
        
        self.soup = BeautifulSoup(self.html_content, 'html.parser')
        
        # 提取完整 CSS 样式
        self.css_content = self._extract_css()
        
        print(f"📄 源文件: {self.html_path.name}")
        print(f"📐 原始尺寸: {self.HTML_WIDTH} × {self.HTML_HEIGHT}")
        print(f"🎯 分辨率: {self.resolution.upper()} ({self.res_desc})")
        print(f"🎨 渲染配置: viewport={self.viewport_w}x{self.viewport_h}, DPI={self.scale_factor}x")
        print(f"📊 PPT 尺寸: {self.prs.slide_width.inches:.3f}\" × {self.prs.slide_height.inches:.3f}\" (16:9)")
    
    def _extract_css(self) -> str:
        """提取所有 CSS 样式"""
        css_parts = []
        for style in self.soup.find_all('style'):
            css_parts.append(style.get_text())
        return '\n'.join(css_parts)
    
    def _get_slide_title(self, slide_elem) -> str:
        """获取幻灯片标题"""
        h1 = slide_elem.find('h1')
        if h1:
            return h1.get_text(strip=True)[:40]
        return "Untitled"
    
    def _create_render_html(self, slide_elem, slide_num: int) -> str:
        """创建用于渲染的独立 HTML"""
        html_template = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <style>
        /* 原始 CSS 样式 */
        {self.css_content}
        
        /* 渲染优化样式 */
        * {{
            -webkit-font-smoothing: antialiased;
            -moz-osx-font-smoothing: grayscale;
            text-rendering: optimizeLegibility;
        }}
        
        /* 隐藏滚动条 */
        *::-webkit-scrollbar {{
            display: none !important;
        }}
        * {{
            scrollbar-width: none !important;
            -ms-overflow-style: none !important;
        }}
        
        /* 重置 body */
        html, body {{
            margin: 0 !important;
            padding: 0 !important;
            background: white !important;
            overflow: hidden !important;
            width: {self.viewport_w}px !important;
            height: {self.viewport_h}px !important;
        }}
        
        /* 强制 slide 样式 - 支持 viewport 缩放 */
        .slide {{
            margin: 0 !important;
            box-shadow: none !important;
            border-radius: 0 !important;
            width: {self.HTML_WIDTH}px !important;
            height: {self.HTML_HEIGHT}px !important;
            position: absolute !important;
            top: 0 !important;
            left: 0 !important;
            transform: scale({self.css_scale}) !important;
            transform-origin: top left !important;
        }}
    </style>
</head>
<body>
    {str(slide_elem)}
</body>
</html>"""
        
        # 保存临时 HTML
        html_file = os.path.join(self.temp_dir, f"slide_{slide_num}.html")
        with open(html_file, 'w', encoding='utf-8') as f:
            f.write(html_template)
        
        return html_file
    
    async def _render_slide_async(self, slide_elem, slide_num: int) -> str:
        """使用 Playwright 渲染单个幻灯片"""
        
        html_file = self._create_render_html(slide_elem, slide_num)
        img_path = os.path.join(self.temp_dir, f"slide_{slide_num}.png")
        
        async with async_playwright() as p:
            # 启动浏览器
            browser = await p.chromium.launch(
                headless=True,
                args=[
                    '--no-sandbox',
                    '--disable-setuid-sandbox',
                    '--disable-dev-shm-usage',
                    '--font-render-hinting=none',  # 更清晰的字体渲染
                ]
            )
            
            # 创建上下文，设置高 DPI
            context = await browser.new_context(
                viewport={
                    'width': int(self.viewport_w),
                    'height': int(self.viewport_h)
                },
                device_scale_factor=self.scale_factor
            )
            
            page = await context.new_page()
            
            # 加载页面
            await page.goto(f'file:///{html_file}', wait_until='networkidle')
            
            # 等待渲染完成
            await page.wait_for_timeout(500)  # 确保所有样式加载完成
            
            # 截图
            await page.screenshot(
                path=img_path,
                full_page=False,
                type='png'
            )
            
            await browser.close()
        
        # 验证截图
        if os.path.exists(img_path):
            img = Image.open(img_path)
            w, h = img.size
            img.close()
            return img_path, w, h
        
        return None, 0, 0
    
    def _render_slide(self, slide_elem, slide_num: int) -> tuple:
        """同步渲染包装"""
        return asyncio.run(self._render_slide_async(slide_elem, slide_num))
    
    def _add_image_to_slide(self, img_path: str) -> bool:
        """将图片添加到 PPT 幻灯片"""
        # 添加空白幻灯片
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置白色背景
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        try:
            # 全屏添加图片
            slide.shapes.add_picture(
                img_path,
                Inches(0), Inches(0),
                self.prs.slide_width,
                self.prs.slide_height
            )
            return True
        except Exception as e:
            print(f"    ❌ 添加图片失败: {e}")
            return False
    
    def convert(self) -> str:
        """执行转换"""
        print(f"\n{'='*60}")
        print(f"🚀 开始转换 html → PPT")
        print(f"{'='*60}\n")
        
        # 查找所有幻灯片
        slides = self.soup.find_all('div', class_='slide')
        total_slides = len(slides)
        
        print(f"📊 发现 {total_slides} 个幻灯片\n")
        
        success_count = 0
        
        for i, slide_elem in enumerate(slides):
            slide_num = i + 1
            title = self._get_slide_title(slide_elem)
            
            print(f"[{slide_num}/{total_slides}] {title}")
            print(f"    🎨 渲染中...", end="", flush=True)
            
            # 渲染幻灯片
            img_path, width, height = self._render_slide(slide_elem, slide_num)
            
            if img_path:
                size_kb = os.path.getsize(img_path) / 1024
                print(f" {width}×{height} ({size_kb:.0f}KB)")
                
                # 添加到 PPT
                print(f"    📝 添加到PPT...", end="", flush=True)
                if self._add_image_to_slide(img_path):
                    print(" ✅")
                    success_count += 1
                else:
                    print(" ❌")
            else:
                print(" ❌ 渲染失败")
        
        # 保存 PPT
        self.prs.save(str(self.output_path))
        
        # 输出结果
        ppt_size = os.path.getsize(self.output_path) / (1024 * 1024)
        
        print(f"\n{'='*60}")
        print(f"✅ 转换完成!")
        print(f"   📁 输出文件: {self.output_path}")
        print(f"   📊 幻灯片数: {success_count}/{total_slides}")
        print(f"   💾 文件大小: {ppt_size:.2f} MB")
        print(f"{'='*60}\n")
        
        # 清理临时文件
        self._cleanup()
        
        return str(self.output_path)
    
    def _cleanup(self):
        """清理临时文件"""
        try:
            shutil.rmtree(self.temp_dir)
        except Exception:
            pass


def main():
    """主函数"""
    import argparse
    
    # 构建分辨率帮助文本
    res_help = "\n".join([f"  {k:6s} - {v[3]} ({v[2]}x DPI)" for k, v in AIZoneHTMLToPPT.RESOLUTIONS.items()])
    
    parser = argparse.ArgumentParser(
        description='html 专用 PPT 转换器 - 支持多种分辨率',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=f"""
分辨率选项:
{res_help}
  Nx     - 自定义 N 倍 DPI (如 6x, 10x)

示例用法:
  python html2ppt.py                           # 默认 4K 分辨率
  python html2ppt.py -r 8k                     # 使用 8K 分辨率
  python html2ppt.py -r 16k                    # 使用 16K 超高清
  python html2ppt.py -r 6x                     # 自定义 6x DPI
  python html2ppt.py -i custom.html -r 4k     # 指定输入文件
  python html2ppt.py -o output.pptx -r 8k     # 指定输出文件
        """
    )
    
    parser.add_argument('-i', '--input', 
                        help='输入 HTML 文件路径 (默认: ai_zone_v2.html)',
                        default=None)
    parser.add_argument('-o', '--output',
                        help='输出 PPT 文件路径 (默认: ai_zone_v2.pptx)',
                        default=None)
    parser.add_argument('-r', '--resolution',
                        help='分辨率: 1080p/2k/4k/5k/8k/16k 或自定义如 6x (默认: 4k)',
                        default='4k')
    
    args = parser.parse_args()
    
    try:
        converter = AIZoneHTMLToPPT(
            html_path=args.input,
            output_path=args.output,
            resolution=args.resolution
        )
        converter.convert()
        return 0
    except FileNotFoundError as e:
        print(f"❌ 错误: {e}")
        return 1
    except Exception as e:
        print(f"❌ 转换失败: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == '__main__':
    exit(main())
